"""
File-type → text extractors. Pure functions; accept bytes (so they work for
both filesystem ingestion and upload-from-browser flows).

`extract(content, filename)` is the single dispatcher. Returns text or empty
string. Optional libraries (docx, pptx, openpyxl) are imported lazily so the
module loads even when those aren't installed.
"""
from __future__ import annotations

import json
import re
import xml.etree.ElementTree as ET
from io import BytesIO
from pathlib import Path

from pypdf import PdfReader
from pypdf.errors import PdfReadError

# File-type buckets. Anything not in EXTRACTABLE returns "".
TEXT_LIKE = {".md", ".txt", ".csv", ".py", ".html", ".json", ".yml", ".yaml",
             ".abap", ".sql", ".sh", ".js", ".ts"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
XLSX_EXT = {".xlsx"}
NOTEBOOK_EXT = {".ipynb"}
XML_LIKE = {".bpmn", ".drawio", ".xml", ".svg"}

EXTRACTABLE = (
    TEXT_LIKE | PDF_EXT | DOCX_EXT | PPTX_EXT | XLSX_EXT | NOTEBOOK_EXT | XML_LIKE
)


def extract(content: bytes, filename: str) -> str:
    """Dispatch to the right extractor based on filename extension."""
    ext = Path(filename).suffix.lower()
    if ext in TEXT_LIKE:
        return _read_text(content)
    if ext in PDF_EXT:
        return _read_pdf(content)
    if ext in DOCX_EXT:
        return _read_docx(content)
    if ext in PPTX_EXT:
        return _read_pptx(content)
    if ext in XLSX_EXT:
        return _read_xlsx(content)
    if ext in NOTEBOOK_EXT:
        return _read_notebook(content)
    if ext in XML_LIKE:
        return _read_xml_like(content)
    return ""


# ---- Per-format readers --------------------------------------------------


def _read_text(content: bytes) -> str:
    try:
        return content.decode("utf-8", errors="ignore")
    except Exception:
        return ""


def _read_pdf(content: bytes) -> str:
    try:
        reader = PdfReader(BytesIO(content))
        out = []
        for i, page in enumerate(reader.pages):
            try:
                t = (page.extract_text() or "").strip()
            except Exception:
                t = ""
            if t:
                out.append(f"[page {i + 1}]\n{t}")
        return "\n\n".join(out)
    except (PdfReadError, Exception):
        return ""


def _read_docx(content: bytes) -> str:
    try:
        import docx as _docx  # python-docx
    except ImportError:
        return ""
    try:
        d = _docx.Document(BytesIO(content))
        parts = [para.text for para in d.paragraphs if para.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""


def _read_pptx(content: bytes) -> str:
    try:
        import pptx as _pptx
    except ImportError:
        return ""
    try:
        prs = _pptx.Presentation(BytesIO(content))
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"[slide {i + 1}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
        return "\n".join(parts)
    except Exception:
        return ""


def _read_xlsx(content: bytes) -> str:
    try:
        from openpyxl import load_workbook as _load_xlsx
    except ImportError:
        return ""
    try:
        wb = _load_xlsx(BytesIO(content), data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        return ""


def _read_notebook(content: bytes) -> str:
    try:
        nb = json.loads(content.decode("utf-8", errors="ignore"))
        parts = []
        for cell in nb.get("cells", []):
            ctype = cell.get("cell_type")
            src = cell.get("source", [])
            if isinstance(src, list):
                src = "".join(src)
            src = src.strip()
            if not src:
                continue
            if ctype == "markdown":
                parts.append(src)
            elif ctype == "code":
                parts.append(f"```python\n{src}\n```")
        return "\n\n".join(parts)
    except Exception:
        return ""


def _read_xml_like(content: bytes) -> str:
    """Pull every text node + every name/label attribute from XML/BPMN/drawio/SVG."""
    try:
        text = content.decode("utf-8", errors="ignore")
        # Strip namespaces for easier traversal
        text_no_ns = re.sub(r"\sxmlns(:\w+)?=\"[^\"]*\"", "", text)
        try:
            root = ET.fromstring(text_no_ns)
        except ET.ParseError:
            # Fall back to plain regex extraction of name= attributes
            names = re.findall(r'\bname=["\'](.*?)["\']', text)
            return "\n".join(n for n in names if n.strip())
        parts = []
        for el in root.iter():
            for attr in ("name", "label", "value"):
                v = el.attrib.get(attr)
                if v and v.strip():
                    parts.append(v.strip())
            if el.text and el.text.strip():
                parts.append(el.text.strip())
        # Dedup adjacent identical lines
        out, prev = [], None
        for line in parts:
            if line != prev:
                out.append(line)
                prev = line
        return "\n".join(out)
    except Exception:
        return ""


def extract_from_path(path: Path) -> str:
    """Convenience for filesystem usage — read bytes then extract."""
    try:
        return extract(path.read_bytes(), path.name)
    except Exception:
        return ""
