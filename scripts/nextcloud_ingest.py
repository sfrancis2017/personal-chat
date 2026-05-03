#!/usr/bin/env python3
"""
Walk a Nextcloud files directory, extract text per file type, chunk + embed
via OpenAI text-embedding-3-small, and insert into the pgvector chunks
table with topic + source_path metadata.

Also runs an authorship + publish-worthiness classifier per file and appends
one row per file to nextcloud-audit.md alongside the script (gitignored —
contains private file titles).

Idempotent: re-running deletes prior chunks for each file path before
re-inserting.

Required env vars (set in /root/ingest/retrieve.env on the droplet):
    OPENAI_API_KEY      = sk-...
    DATABASE_URL        = postgresql://...
    NEXTCLOUD_ROOT      = absolute path to Nextcloud files directory
    AUTHOR_SIGNALS      = pipe-separated regex of author identifiers (name, handles, projects)
    OWNER_PREFIXES      = comma-separated filename prefixes that mark owner content (e.g. "owner_,mine_")
"""
import json
import os
import re
import sys
import time
import xml.etree.ElementTree as ET
from pathlib import Path

import psycopg2
import tiktoken
from openai import OpenAI
from pgvector.psycopg2 import register_vector
from pypdf import PdfReader
from pypdf.errors import PdfReadError

# OpenAI text-embedding-3-small uses cl100k_base; max 8192 tokens / call.
ENC = tiktoken.get_encoding("cl100k_base")
MAX_EMBED_TOKENS = 8000  # conservative ceiling (limit is 8192)

# Optional extractors — gracefully skip the file type if not installed
try:
    import docx as _docx
except ImportError:
    _docx = None
try:
    import pptx as _pptx
except ImportError:
    _pptx = None
try:
    from openpyxl import load_workbook as _load_xlsx
except ImportError:
    _load_xlsx = None

NEXTCLOUD_ROOT = Path(os.environ["NEXTCLOUD_ROOT"])  # required — set in retrieve.env
AUDIT_PATH = Path(__file__).parent / "nextcloud-audit.md"

CHUNK_TOKENS = 1000      # ~750 words; well under the 8192-token embed cap
CHUNK_OVERLAP_TOK = 100  # token overlap between adjacent chunks

# File types we extract text from. Anything else is skipped with a count.
TEXT_LIKE = {".md", ".txt", ".csv", ".py", ".html", ".json", ".yml", ".yaml",
             ".abap", ".sql", ".sh", ".js", ".ts"}
PDF_EXT = {".pdf"}
DOCX_EXT = {".docx"}
PPTX_EXT = {".pptx"}
XLSX_EXT = {".xlsx"}
NOTEBOOK_EXT = {".ipynb"}
XML_LIKE = {".bpmn", ".drawio", ".xml", ".svg"}

EXTRACTABLE = (
    TEXT_LIKE | PDF_EXT | DOCX_EXT | PPTX_EXT | XLSX_EXT | NOTEBOOK_EXT | XML_LIKE
)

# Skip patterns
SKIP_DIR_NAMES = {".attachments.9974", "Templates"}  # Nextcloud system / boilerplate
SKIP_FILE_PATTERNS = [
    re.compile(r"^\.(tmp|trash|version)"),  # hidden Nextcloud metadata
]

# ---- Extractors ----------------------------------------------------------


def read_text(p: Path) -> str:
    try:
        return p.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def read_pdf(p: Path) -> str:
    try:
        reader = PdfReader(str(p))
        out = []
        for i, page in enumerate(reader.pages):
            try:
                t = (page.extract_text() or "").strip()
            except Exception:
                t = ""
            if t:
                out.append(f"[page {i + 1}]\n{t}")
        return "\n\n".join(out)
    except (PdfReadError, Exception):
        return ""


def read_docx(p: Path) -> str:
    if _docx is None:
        return ""
    try:
        d = _docx.Document(str(p))
        parts = [para.text for para in d.paragraphs if para.text.strip()]
        for tbl in d.tables:
            for row in tbl.rows:
                parts.append(" | ".join(cell.text for cell in row.cells))
        return "\n".join(parts)
    except Exception:
        return ""


def read_pptx(p: Path) -> str:
    if _pptx is None:
        return ""
    try:
        prs = _pptx.Presentation(str(p))
        parts = []
        for i, slide in enumerate(prs.slides):
            parts.append(f"[slide {i + 1}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
        return "\n".join(parts)
    except Exception:
        return ""


def read_xlsx(p: Path) -> str:
    if _load_xlsx is None:
        return ""
    try:
        wb = _load_xlsx(str(p), data_only=True, read_only=True)
        parts = []
        for sheet in wb.worksheets:
            parts.append(f"[sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception:
        return ""


def read_notebook(p: Path) -> str:
    try:
        nb = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
        parts = []
        for cell in nb.get("cells", []):
            ctype = cell.get("cell_type")
            src = cell.get("source", [])
            if isinstance(src, list):
                src = "".join(src)
            src = src.strip()
            if not src:
                continue
            if ctype == "markdown":
                parts.append(src)
            elif ctype == "code":
                parts.append(f"```python\n{src}\n```")
        return "\n\n".join(parts)
    except Exception:
        return ""


def read_xml_like(p: Path) -> str:
    """Pull every text node + every name/label attribute from XML/BPMN/drawio/SVG."""
    try:
        text = p.read_text(encoding="utf-8", errors="ignore")
        # Strip namespaces for easier traversal
        text_no_ns = re.sub(r"\sxmlns(:\w+)?=\"[^\"]*\"", "", text)
        try:
            root = ET.fromstring(text_no_ns)
        except ET.ParseError:
            # Fall back to plain regex extraction of name= attributes + text content
            names = re.findall(r'\bname=["\'](.*?)["\']', text)
            return "\n".join(n for n in names if n.strip())
        parts = []
        for el in root.iter():
            for attr in ("name", "label", "value"):
                v = el.attrib.get(attr)
                if v and v.strip():
                    parts.append(v.strip())
            if el.text and el.text.strip():
                parts.append(el.text.strip())
        # Dedup adjacent identical lines
        out, prev = [], None
        for line in parts:
            if line != prev:
                out.append(line)
                prev = line
        return "\n".join(out)
    except Exception:
        return ""


def extract(p: Path) -> str:
    ext = p.suffix.lower()
    if ext in TEXT_LIKE:
        return read_text(p)
    if ext in PDF_EXT:
        return read_pdf(p)
    if ext in DOCX_EXT:
        return read_docx(p)
    if ext in PPTX_EXT:
        return read_pptx(p)
    if ext in XLSX_EXT:
        return read_xlsx(p)
    if ext in NOTEBOOK_EXT:
        return read_notebook(p)
    if ext in XML_LIKE:
        return read_xml_like(p)
    return ""


# ---- Topic + chunking ----------------------------------------------------


def topic_of(rel_path: Path) -> str:
    parts = rel_path.parts
    # Files at the root (no folder) bucket under "nextcloud-root" so each
    # loose file doesn't become its own topic.
    if len(parts) <= 1:
        return "nextcloud-root"
    top = parts[0]
    return (
        re.sub(r"[^a-zA-Z0-9]+", "-", top).strip("-").lower() or "nextcloud-misc"
    )


def clean_text(text: str) -> str:
    """Strip NUL bytes (Postgres rejects them) and other non-printable junk."""
    if not text:
        return ""
    # Postgres TEXT can't store NUL; strip explicitly
    text = text.replace("\x00", "")
    # Replace stray carriage returns; collapse other control chars except \n \t
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    return text


def chunk_text(text: str):
    """Token-based chunking. Guarantees each chunk is under the embed API limit."""
    text = clean_text(text)
    if not text:
        return []
    tokens = ENC.encode(text, disallowed_special=())
    out = []
    i = 0
    step = max(1, CHUNK_TOKENS - CHUNK_OVERLAP_TOK)
    while i < len(tokens):
        slice_tokens = tokens[i:i + CHUNK_TOKENS]
        chunk = ENC.decode(slice_tokens)
        # Defensive truncation if decode somehow produced too-long text
        if len(ENC.encode(chunk, disallowed_special=())) > MAX_EMBED_TOKENS:
            chunk = ENC.decode(slice_tokens[:MAX_EMBED_TOKENS])
        chunk = chunk.strip()
        if chunk:
            out.append(chunk)
        i += step
    return out


# ---- Authorship + publish-worthiness classifier --------------------------

PUBLISHER_BOILERPLATE = re.compile(
    r"\b(ISBN[\s:-]*\d|©\s*\d{4}\s*(SAP|Pearson|O'?Reilly|McGraw|Wiley|Elsevier|"
    r"Springer|Cambridge|Oxford|Manning|Apress|Packt|Addison)|"
    r"All rights reserved\.|"
    r"SAP\s+Press|"
    r"This document is the property of)\b",
    re.IGNORECASE,
)

# Owner-identifying tokens — set via env to avoid baking personal handles
# into committed source. Default is a no-op pattern that won't match anything.
_AUTHOR_PATTERN = os.environ.get("AUTHOR_SIGNALS", "__no_match__")
AUTHOR_SIGNALS = re.compile(rf"\b({_AUTHOR_PATTERN})\b", re.IGNORECASE)

# Filename prefixes that mark owner-authored content (e.g. "owner_,mine_").
OWNER_PREFIXES = tuple(
    p.strip().lower()
    for p in os.environ.get("OWNER_PREFIXES", "").split(",")
    if p.strip()
)


def classify_authorship(rel_path: Path, text: str) -> tuple[str, str]:
    """Return (author_confidence, reason). Conservative — defaults to UNCERTAIN."""
    p = str(rel_path).lower()

    # Strong third-party folders
    if "bpmn - s4hana" in p:
        return "high_third_party", "SAP Best Practices process content"
    if "sap icons" in p:
        return "high_third_party", "SAP icon assets"
    if "sap press" in p:
        return "high_third_party", "SAP Press materials"

    # Owner-prefixed filenames
    if OWNER_PREFIXES and any(rel_path.name.lower().startswith(pre) for pre in OWNER_PREFIXES):
        return "high_owner", "explicit owner prefix"

    # Boilerplate scan
    head = text[:5000]
    if PUBLISHER_BOILERPLATE.search(head):
        return "high_third_party", "publisher boilerplate detected in opening"

    # Owner-identifying tokens in content
    owner_hits = len(AUTHOR_SIGNALS.findall(text[:20000]))
    if owner_hits >= 2:
        return "high_owner", f"{owner_hits} self-references in content"

    return "uncertain", "no decisive signal — needs review"


def publish_worthiness(rel_path: Path, text: str) -> tuple[int, str]:
    """1-5 score. Heuristic only; the audit is for human review, not auto-publish."""
    n = len(text)
    if n < 500:
        return 1, "too short"
    if n < 3000:
        return 2, "short"
    if n < 15000:
        score = 3
    elif n < 50000:
        score = 4
    else:
        score = 5
    name = rel_path.name.lower()
    if any(s in name for s in ("draft", "wip", "todo", "tmp")):
        score = max(1, score - 2)
    return score, f"{n} chars"


def suggested_folder(rel_path: Path) -> str:
    p = str(rel_path).lower()
    if "bpmn" in p or "process" in p:
        return "architecture/business-process/"
    if "diagram" in p:
        return "architecture/"
    if "agent" in p or "llm" in p or "rag" in p:
        return "ai/agents-and-tools/"
    if "sap" in p:
        return "reference/sap/"
    if rel_path.suffix.lower() == ".py":
        return "software-engineering/"
    return "architecture/"


# ---- Embedding + DB ------------------------------------------------------


def embed(client: OpenAI, text: str):
    r = client.embeddings.create(model="text-embedding-3-small", input=text)
    return r.data[0].embedding


# ---- Main walk -----------------------------------------------------------


def should_skip(p: Path) -> bool:
    if any(part in SKIP_DIR_NAMES for part in p.parts):
        return True
    if any(pat.match(p.name) for pat in SKIP_FILE_PATTERNS):
        return True
    return False


def main():
    if not NEXTCLOUD_ROOT.exists():
        print(f"ERROR: {NEXTCLOUD_ROOT} not found", file=sys.stderr)
        sys.exit(1)

    client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    register_vector(conn)
    cur = conn.cursor()

    files = []
    for p in NEXTCLOUD_ROOT.rglob("*"):
        if not p.is_file():
            continue
        if should_skip(p):
            continue
        if p.suffix.lower() not in EXTRACTABLE:
            continue
        files.append(p)

    print(f"Found {len(files)} extractable files under {NEXTCLOUD_ROOT}")

    # Open audit file (overwrite each run)
    audit = AUDIT_PATH.open("w", encoding="utf-8")
    audit.write("# Nextcloud audit\n\n")
    audit.write("Generated by `nextcloud_ingest.py`. Each row is one file from the "
                "Nextcloud files tree, with the classifier's verdict on authorship "
                "and publish-worthiness.\n\n")
    audit.write("**Hard rule for publishing**: only rows with `author_confidence = high_owner` "
                "AND `publish_worthy >= 3` should be considered. UNCERTAIN rows get "
                "manual review. high_third_party rows are NEVER published.\n\n")
    audit.write("| File | Author | Reason | Worth | Suggested folder | Chunks |\n")
    audit.write("|---|---|---|---|---|---|\n")

    total_chunks = 0
    skipped = 0
    for i, p in enumerate(files, 1):
        rel = p.relative_to(NEXTCLOUD_ROOT)
        try:
            text = extract(p)
        except Exception as e:
            print(f"[{i}/{len(files)}] {rel}  EXTRACT ERR {e}")
            skipped += 1
            continue

        if not text or len(text) < 50:
            print(f"[{i}/{len(files)}] {rel}  skip (empty)")
            skipped += 1
            continue

        topic = topic_of(rel)
        author, author_reason = classify_authorship(rel, text)
        worth, worth_reason = publish_worthiness(rel, text)
        sugg = suggested_folder(rel)

        chunks = chunk_text(text)
        # Re-ingest: clear prior rows for this path
        cur.execute("DELETE FROM chunks WHERE source_path = %s", (str(rel),))

        ok = 0
        for ch in chunks:
            try:
                emb = embed(client, ch)
            except Exception as e:
                print(f"  embed err: {e}")
                continue
            cur.execute(
                "INSERT INTO chunks (source_url, source_path, text, embedding, metadata, topic) "
                "VALUES (%s, %s, %s, %s, %s, %s)",
                (
                    None,
                    str(rel),
                    ch,
                    emb,
                    json.dumps(
                        {
                            "origin": "nextcloud",
                            "title": rel.name,
                            "author_confidence": author,
                            "publish_worthy": worth,
                        }
                    ),
                    topic,
                ),
            )
            ok += 1
        conn.commit()
        total_chunks += ok

        audit_file_link = str(rel).replace("|", "\\|")
        audit.write(
            f"| `{audit_file_link}` | {author} | {author_reason} | "
            f"{worth} ({worth_reason}) | `{sugg}` | {ok} |\n"
        )
        audit.flush()

        print(
            f"[{i}/{len(files)}] {rel}  "
            f"chunks={ok} topic={topic} author={author}"
        )
        time.sleep(0.05)  # gentle rate limit

    audit.close()
    cur.close()
    conn.close()

    print()
    print(f"Done. {total_chunks} chunks inserted from {len(files) - skipped} files "
          f"({skipped} skipped/empty).")
    print(f"Audit written to {AUDIT_PATH}")


if __name__ == "__main__":
    main()
